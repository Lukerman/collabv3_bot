"""
Text extraction utilities for CollaLearn bot.
Extracts text from various file formats.
"""

import logging
import io
from typing import Optional

logger = logging.getLogger(__name__)


async def extract_text_from_file(
    file_bytes: bytes,
    file_name: str,
    mime_type: str
) -> Optional[str]:
    """
    Extract text content from various file types.
    
    Args:
        file_bytes: File content as bytes
        file_name: Name of the file
        mime_type: MIME type of the file
    
    Returns:
        Extracted text or None if extraction fails
    """
    try:
        # PDF files
        if mime_type == "application/pdf":
            return await extract_from_pdf(file_bytes)
        
        # DOCX files
        elif "wordprocessingml" in mime_type:
            return await extract_from_docx(file_bytes)
        
        # PPTX files
        elif "presentationml" in mime_type:
            return await extract_from_pptx(file_bytes)
        
        # Text files
        elif mime_type == "text/plain":
            return file_bytes.decode('utf-8', errors='ignore')
        
        # Images (OCR would be needed - placeholder)
        elif mime_type.startswith("image/"):
            logger.info(f"Image file detected: {file_name}. OCR not implemented.")
            return None
        
        else:
            logger.warning(f"Unsupported file type for text extraction: {mime_type}")
            return None
            
    except Exception as e:
        logger.error(f"Error extracting text from {file_name}: {e}")
        return None


async def extract_from_pdf(file_bytes: bytes) -> Optional[str]:
    """
    Extract text from PDF file.
    
    Args:
        file_bytes: PDF file content
    
    Returns:
        Extracted text or None
    """
    try:
        import PyPDF2
        
        pdf_file = io.BytesIO(file_bytes)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        
        return text.strip()
    
    except ImportError:
        logger.error("PyPDF2 not installed. Install with: pip install PyPDF2")
        return None
    except Exception as e:
        logger.error(f"Error extracting from PDF: {e}")
        return None


async def extract_from_docx(file_bytes: bytes) -> Optional[str]:
    """
    Extract text from DOCX file.
    
    Args:
        file_bytes: DOCX file content
    
    Returns:
        Extracted text or None
    """
    try:
        import docx
        
        doc_file = io.BytesIO(file_bytes)
        doc = docx.Document(doc_file)
        
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        return text.strip()
    
    except ImportError:
        logger.error("python-docx not installed. Install with: pip install python-docx")
        return None
    except Exception as e:
        logger.error(f"Error extracting from DOCX: {e}")
        return None


async def extract_from_pptx(file_bytes: bytes) -> Optional[str]:
    """
    Extract text from PPTX file.
    
    Args:
        file_bytes: PPTX file content
    
    Returns:
        Extracted text or None
    """
    try:
        from pptx import Presentation
        
        pptx_file = io.BytesIO(file_bytes)
        prs = Presentation(pptx_file)
        
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        return text.strip()
    
    except ImportError:
        logger.error("python-pptx not installed. Install with: pip install python-pptx")
        return None
    except Exception as e:
        logger.error(f"Error extracting from PPTX: {e}")
        return None


def truncate_text(text: str, max_length: int) -> str:
    """
    Truncate text to maximum length.
    
    Args:
        text: Text to truncate
        max_length: Maximum length
    
    Returns:
        Truncated text
    """
    if len(text) <= max_length:
        return text
    
    return text[:max_length] + "..."